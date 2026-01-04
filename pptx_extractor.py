
"""
pptx_extractor.py
Usage: python pptx_extractor.py input.pptx out_dir

What it does:
- Extracts text (paragraphs + runs) and run-level font properties
- Extracts images (pictures) to out_dir/images/
- Extracts tables to CSV files
- Extracts slide notes
- Produces per-slide JSON metadata with shape bbox (emu) and normalized coords
- Flags "problematic" shapes (SmartArt, Chart, GraphicFrame w/o accessible image)
- Saves XML snapshot for flagged shapes for debugging
- Produces per-slide bbox CSV for quick use
"""

import sys
import json
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from lxml import etree
from PIL import Image
import io
import csv

EMU_PER_INCH = 914400  # pptx EMU constant
DPI = 96  # used only if you need px conversion; optional

def emu_to_inches(emu):
    return emu / EMU_PER_INCH

def emu_to_px(emu, dpi=DPI):
    inches = emu_to_inches(emu)
    return int(round(inches * dpi))

def norm(val, total):
    return float(val) / float(total) if total else 0.0

def shape_xml_string(shape):
    # Return pretty printed XML string of the shape.element for debugging
    el = shape.element
    return etree.tostring(el, pretty_print=True, encoding='unicode')

def ensure_dir(p: Path):
    p.mkdir(parents=True, exist_ok=True)

def save_image_blob(blob_bytes: bytes, outpath: Path):
    try:
        outpath.write_bytes(blob_bytes)
        return True
    except Exception as e:
        # Try via PIL for some compatibility if needed
        try:
            im = Image.open(io.BytesIO(blob_bytes))
            im.save(outpath)
            return True
        except Exception:
            return False

def extract(pptx_path: Path, out_dir: Path):
    ensure_dir(out_dir)
    images_dir = out_dir / "images"
    ensure_dir(images_dir)
    metadata = {"source": str(pptx_path), "slides": []}

    prs = Presentation(str(pptx_path))
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    for si, slide in enumerate(prs.slides, start=1):
        slide_meta = {
            "slide_index": si - 1,
            "slide_num": si,
            "width_emu": slide_width,
            "height_emu": slide_height,
            "shapes": [],
            "notes": None,
            "problem_shapes": []
        }
        slide_dir = out_dir / f"slide_{si:02d}"
        ensure_dir(slide_dir)

        # Save notes if present
        notes_text = ""
        try:
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                notes_text = "\n".join([p.text for p in notes_slide.notes_text_frame.paragraphs])
                (slide_dir / "notes.txt").write_text(notes_text, encoding="utf-8")
                slide_meta["notes"] = notes_text
        except Exception:
            # Some slides may not have notes_slide attribute
            slide_meta["notes"] = None

        # Iterate shapes (including groups)
        def handle_shape(shape, parent_index=None, depth=0):
            sidx = len(slide_meta["shapes"]) + 1
            base = {
                "shape_index": sidx,
                "type": str(shape.shape_type),
                "is_group": shape.shape_type == MSO_SHAPE_TYPE.GROUP,
                "left_emu": getattr(shape, "left", None),
                "top_emu": getattr(shape, "top", None),
                "width_emu": getattr(shape, "width", None),
                "height_emu": getattr(shape, "height", None),
                "left_norm": None,
                "top_norm": None,
                "width_norm": None,
                "height_norm": None,
                "has_text": False,
                "text": None,
                "paragraphs": [],
                "has_table": False,
                "table_rows": None,
                "table_cols": None,
                "has_image": False,
                "image_path": None,
                "notes": None,
                "xml_snapshot": None
            }

            # normalize bbox
            try:
                l = base["left_emu"] or 0
                t = base["top_emu"] or 0
                w = base["width_emu"] or 0
                h = base["height_emu"] or 0
                base["left_norm"] = norm(l, slide_width)
                base["top_norm"] = norm(t, slide_height)
                base["width_norm"] = norm(w, slide_width)
                base["height_norm"] = norm(h, slide_height)
            except Exception:
                pass

            # Text extraction
            try:
                if getattr(shape, "has_text_frame", False):
                    base["has_text"] = True
                    base["text"] = shape.text  # flattened text
                    paras = []
                    for p in shape.text_frame.paragraphs:
                        runs = []
                        for r in p.runs:
                            font = r.font
                            run_meta = {
                                "text": r.text,
                                "bold": bool(font.bold) if font is not None else None,
                                "italic": bool(font.italic) if font is not None else None,
                                "underline": bool(font.underline) if font is not None else None,
                                "font_name": font.name if font is not None else None,
                                "font_size_pt": font.size.pt if (font is not None and font.size is not None) else None,
                                "color_rgb": None
                            }
                            # color may be complicated; best-effort
                            try:
                                if font.color and font.color.rgb:
                                    run_meta["color_rgb"] = font.color.rgb
                            except Exception:
                                pass
                            runs.append(run_meta)
                        is_bullet = False
                        try:
                            if p.level is not None:
                                # python-pptx marks bullet via `p._p.pPr.bu*`
                                if p._p.pPr is not None and (
                                    p._p.pPr.find(".//a:buChar", namespaces=p._p.nsmap) is not None or
                                    p._p.pPr.find(".//a:buAutoNum", namespaces=p._p.nsmap) is not None
                                ):
                                    is_bullet = True
                        except Exception:
                            pass
                            
                        paras.append({
                            "level": p.level,
                            "is_bullet": is_bullet,
                            "runs": runs
                        })

                    base["paragraphs"] = paras
            except Exception:
                # fallback: include XML snapshot to debug
                base["xml_snapshot"] = shape_xml_string(shape)

            # Table extraction
            try:
                if getattr(shape, "has_table", False):
                    tbl = shape.table
                    base["has_table"] = True
                    base["table_rows"] = len(tbl.rows)
                    base["table_cols"] = len(tbl.columns)
                    # dump CSV
                    csv_path = slide_dir / f"shape_{sidx}_table.csv"
                    with csv_path.open("w", encoding="utf-8", newline="") as cf:
                        writer = csv.writer(cf)
                        for r in range(len(tbl.rows)):
                            rowvals = []
                            for c in range(len(tbl.columns)):
                                rowvals.append(tbl.cell(r, c).text.replace("\n", " ").strip())
                            writer.writerow(rowvals)
                    base["table_csv"] = str(csv_path.relative_to(out_dir))
            except Exception:
                base["has_table"] = False

            # Image extraction
            try:
                # Preferred: shape.image (works for Picture shapes)
                img = getattr(shape, "image", None)
                if img is not None:
                    base["has_image"] = True
                    ext = img.ext
                    blob = img.blob
                    img_name = f"slide{si:02d}_shape{sidx}_img.{ext}"
                    img_path = images_dir / img_name
                    saved = save_image_blob(blob, img_path)
                    if saved:
                        base["image_path"] = str(img_path.relative_to(out_dir))
            except Exception:
                # for grouped images or complex cases, fallback to xml snapshot flag
                pass

            # Flag problematic shapes: charts, smartart, graphicFrame w/o accessible image
            try:
                tag = shape.element.tag.lower()
                # Many charts/smartart are graphicFrame (a:graphicFrame). Decide by checking for c:chart or <p:graphicFrame> with chart namespace
                xml_str = shape_xml_string(shape)
                # heuristics
                if ("chart" in xml_str.lower()) or ("graphicdata" in xml_str.lower() and "chart" in xml_str.lower()) or ("smartart" in xml_str.lower()):
                    # mark as problem (we may want to export as image instead)
                    base["xml_snapshot"] = xml_str
                    base["problem"] = "chart_or_smartart"
                    slide_meta["problem_shapes"].append(sidx)
                # a graphicFrame that doesn't have an accessible image: flag
                elif shape.shape_type == MSO_SHAPE_TYPE.GRAPHICAL_FRAME or ("graphicFrame" in xml_str and base["has_image"] is False):
                    # check if there's an embedded image (blip)
                    if "blip" not in xml_str.lower():
                        base["xml_snapshot"] = xml_str
                        base["problem"] = "graphic_frame_no_blip"
                        slide_meta["problem_shapes"].append(sidx)
            except Exception:
                # ignore
                pass

            slide_meta["shapes"].append(base)

            # If group, recurse
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                try:
                    for sub in shape.shapes:
                        handle_shape(sub, parent_index=sidx, depth=depth+1)
                except Exception:
                    pass

        # iterate top-level shapes
        for shape in slide.shapes:
            handle_shape(shape)

        # Save per-slide bbox CSV
        bbox_csv = slide_dir / f"slide_{si:02d}_bboxes.csv"
        with bbox_csv.open("w", encoding="utf-8", newline="") as bf:
            writer = csv.writer(bf)
            writer.writerow(["shape_index","left_emu","top_emu","width_emu","height_emu","left_norm","top_norm","width_norm","height_norm","has_text","has_table","has_image","image_path","problem"])
            for s in slide_meta["shapes"]:
                writer.writerow([
                    s.get("shape_index"),
                    s.get("left_emu"),
                    s.get("top_emu"),
                    s.get("width_emu"),
                    s.get("height_emu"),
                    s.get("left_norm"),
                    s.get("top_norm"),
                    s.get("width_norm"),
                    s.get("height_norm"),
                    s.get("has_text"),
                    s.get("has_table"),
                    s.get("has_image"),
                    s.get("image_path"),
                    s.get("problem", "")
                ])

        # Save any XML snapshots for problematic shapes
        if slide_meta["problem_shapes"]:
            ps_dir = slide_dir / "problem_shapes_xml"
            ensure_dir(ps_dir)
            for s in slide_meta["shapes"]:
                if s.get("problem"):
                    idx = s["shape_index"]
                    fname = ps_dir / f"shape_{idx}_xml.xml"
                    xmltxt = s.get("xml_snapshot", "<no-xml>")
                    fname.write_text(xmltxt, encoding="utf-8")

        # Save slide json
        slide_json_path = slide_dir / f"slide_{si:02d}_metadata.json"
        slide_json_path.write_text(json.dumps(slide_meta, ensure_ascii=False, indent=2), encoding="utf-8")

        metadata["slides"].append({"slide_num": si, "meta_path": str(slide_json_path.relative_to(out_dir))})

    # Save global metadata
    (out_dir / "metadata.json").write_text(json.dumps(metadata, ensure_ascii=False, indent=2), encoding="utf-8")
    print(f"Extraction complete. Output written to: {out_dir}")

if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python pptx_extractor.py input.pptx out_dir")
        sys.exit(1)
    pptx_path = Path(sys.argv[1])
    out_dir = Path(sys.argv[2])
    extract(pptx_path, out_dir)


