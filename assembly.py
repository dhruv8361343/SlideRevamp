from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import csv
from PIL import Image
import os
from pathlib import Path
from pptx.enum.text import MSO_AUTO_SIZE,PP_ALIGN

prs = Presentation()


ASSETS_DIR = Path("/kaggle/working/outputs")


def apply_background(prs, slide, bg_path):
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    slide.shapes.add_picture(
        str(bg_path),
        0,
        0,
        slide_width,
        slide_height
    )


def n2pt(prs, x, y, w, h):
    return (
        int(x * prs.slide_width),
        int(y * prs.slide_height),
        int(w * prs.slide_width),
        int(h * prs.slide_height)
    )
def safe_apply_color(font, color_data):
    """Robustly applies color from various formats (hex, tuple, name)."""
    if not color_data:
        return

    # Clean string
    c_str = str(color_data).lower().replace("#", "").replace(" ", "").replace("(", "").replace(")", "")
    
    # Handle known names manually if needed
    if "green" in c_str: 
        font.color.rgb = RGBColor(0, 128, 0)
        return
        
    # Handle Hex (6 chars)
    if len(c_str) == 6:
        try:
            font.color.rgb = RGBColor.from_string(c_str)
        except: pass
    
    # Handle Tuple string '0,255,0'
    elif "," in c_str:
        try:
            parts = c_str.split(",")
            if len(parts) >= 3:
                font.color.rgb = RGBColor(int(parts[0]), int(parts[1]), int(parts[2]))
        except: pass

def parse_color(color_val):
    """Handles Hex, Tuples, and common named colors."""
    if not color_val: return None
    c = str(color_val).lower().strip().replace("#", "")
    
    # Handle common names
    named_colors = {"green": "008000", "red": "FF0000", "blue": "0000FF", "white": "FFFFFF", "black": "000000"}
    if c in named_colors: return RGBColor.from_string(named_colors[c])
    
    # Handle Hex
    if len(c) == 6:
        try: return RGBColor.from_string(c)
        except: pass
        
    # Handle Tuple strings like "(0, 128, 0)"
    if "," in c:
        try:
            rgb = [int(x.strip()) for x in c.replace("(", "").replace(")", "").split(",")]
            return RGBColor(rgb[0], rgb[1], rgb[2])
        except: pass
    return None
                
def apply_run_color(font, run_data, el):
    # Order of priority: 1. run color_rgb, 2. run font_color, 3. element font_color
    val = None
    if isinstance(run_data, dict):
        val = run_data.get("color_rgb") or run_data.get("font_color")
    
    # If run didn't have it, check the parent element
    if not val:
        val = el.get("font_color")

    if not val:
        return

    try:
        if hasattr(val, 'rgb'):
            font.color.rgb = val
            return

        if isinstance(val, int):
            r, g, b = (val >> 16) & 0xFF, (val >> 8) & 0xFF, val & 0xFF
            font.color.rgb = RGBColor(r, g, b)
            return

        clean = str(val).replace("#", "").strip()
        if len(clean) == 6:
            font.color.rgb = RGBColor.from_string(clean)

    except Exception:
        font.color.rgb = RGBColor(0, 0, 0) # Fallback to black

def apply_font_color(font, color_val):
    """
    Robustly applies color to a font object.
    Supports Hex strings (from updated extractor), integers, and basic names.
    """
    if not color_val:
        return
    
    try:
        # 1. Handle Hex strings (e.g., "FF5733" or "#FF5733")
        if isinstance(color_val, str):
            clean_hex = color_val.strip().replace("#", "").lower()
            
            # Map common names if the extractor picked them up
            named_colors = {
                "green": "008000", "red": "ff0000", "blue": "0000ff", 
                "white": "ffffff", "black": "000000"
            }
            hex_to_use = named_colors.get(clean_hex, clean_hex)
            
            if len(hex_to_use) == 6:
                font.color.rgb = RGBColor.from_string(hex_to_use)
                
        # 2. Handle Integers (fallback for old extractor versions)
        elif isinstance(color_val, int):
            r = (color_val >> 16) & 0xFF
            g = (color_val >> 8) & 0xFF
            b = color_val & 0xFF
            font.color.rgb = RGBColor(r, g, b)
            
    except Exception as e:
        print(f"Warning: Could not apply color '{color_val}': {e}")

def enable_bullets(paragraph):
    """
    Manually enables bullets by injecting the correct XML elements 
    into the paragraph properties.
    """
    pPr = paragraph._p.get_or_add_pPr()
    
    # 1. Check if buChar already exists, if not, create it
    # The XML tag for bullet character is 'buChar'
    buChar = pPr.find('.//a:buChar', namespaces=pPr.nsmap)
    if buChar is None:
        from pptx.oxml.xmlchemy import OxmlElement
        buChar = OxmlElement('a:buChar')
        pPr.append(buChar)
    
    # 2. Set the bullet character (e.g., a standard dot)
    buChar.set('char', '•')




                
from pptx.enum.text import MSO_AUTO_SIZE 

def add_text(prs,slide, el):
    # expand usable content area automatically
    l, t, w, h = n2pt(prs,el["x"], el["y"], el["width"], el["height"])

    box = slide.shapes.add_textbox(l, t, w, h)

    tf = box.text_frame
    tf.word_wrap = True
    
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE 

    content_data = el.get("content", []) 
    alignment_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
    align_enum = alignment_map.get(el.get("align", "left"), PP_ALIGN.LEFT)
    

    for i, para_data in enumerate(content_data):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align_enum

        ls = el.get("line_spacing")
        if ls:
            p.line_spacing = ls

        
        if isinstance(para_data, dict):
            p.level = para_data.get("level", 0)
            if para_data.get("is_bullet"):
                enable_bullets(p)
            runs_to_process = para_data.get("runs", [])
            if para_data.get("is_bullet", False):
                enable_bullets(p)
        else:
            # If para_data is just a list of runs directly
            runs_to_process = para_data
        
        
        for run_data in runs_to_process:
            run = p.add_run()
            if isinstance(run_data, str):
                run.text = run_data
                font = run.font
                apply_font_color(font, el.get("font_color"))

            
            else:
                run.text = run_data.get("text", "")
                font = run.font
                if run_data.get("bold"): font.bold = True
                    
                specific_color = run_data.get("color_rgb") or para_data.get("font_color") if isinstance(para_data, dict) else None
                final_color = specific_color or el.get("font_color")
                
                apply_font_color(font, final_color)


            # Apply font size (either from run, element, or default 18)
            size = 18
            if isinstance(run_data, dict):
                size = run_data.get("font_size") or el.get("font_size", 18)
            else:
                size = el.get("font_size", 18)
            font.size = Pt(size)



def add_image(prs,slide, el):
    l, t, w, h = n2pt(prs,el["x"], el["y"], el["width"], el["height"])
    
    # Construct full path
    image_path = ASSETS_DIR / el["source"]
    
    if not image_path.exists():
        print(f"Image missing: {image_path}")
        return

    # 1. Load to get dimensions
    with Image.open(image_path) as img:
        img_w, img_h = img.size

    frame_ratio = w / h
    img_ratio = img_w / img_h

    if el.get("width", 0) > 0.9 and el.get("height", 0) > 0.9:
        slide.shapes.add_picture(str(image_path), 0, 0, prs.slide_width, prs.slide_height)
        return
    
    # 2. Add Picture
    pic = slide.shapes.add_picture(str(image_path), l, t)
    
    # 3. Check Fit Mode (Calculated in previous step)
    fit_mode = el.get("fit", "contain") # Default to contain for safety

    if fit_mode == "cover":
        # --- CROP TO FILL ---
        if frame_ratio > img_ratio:
            # Frame is wider: Crop Top/Bottom
            target_h = img_w / frame_ratio
            crop_amount = (img_h - target_h) / img_h
            pic.crop_top = crop_amount / 2
            pic.crop_bottom = crop_amount / 2
        else:
            # Frame is taller: Crop Left/Right
            target_w = img_h * frame_ratio
            crop_amount = (img_w - target_w) / img_w
            pic.crop_left = crop_amount / 2
            pic.crop_right = crop_amount / 2
            
        # Resize to fill box
        pic.left, pic.top, pic.width, pic.height = l, t, w, h

    else:
        # --- CONTAIN (Fit Inside) --- 
        # We need to center the image inside the box without stretching
        if frame_ratio > img_ratio:
            # Frame is wider than image: Fit to Height
            new_w = h * img_ratio
            left_offset = (w - new_w) / 2
            pic.height = h
            pic.width = int(new_w)
            pic.left = l + int(left_offset)
            pic.top = t
        else:
            # Frame is taller than image: Fit to Width
            new_h = w / img_ratio
            top_offset = (h - new_h) / 2
            pic.width = w
            pic.height = int(new_h)
            pic.left = l
            pic.top = t + int(top_offset)



def add_table(prs,slide, el):
    csv_path = ASSETS_DIR / el["source"]
    
    if not csv_path.exists():
        return

    with open(csv_path, newline='', encoding="utf-8") as f:
        rows = list(csv.reader(f))

    if not rows: return

    n_rows = len(rows)
    n_cols = max(len(r) for r in rows) # Safety max
    l, t, w, h = n2pt(prs,el["x"], el["y"], el["width"], el["height"])

    table = slide.shapes.add_table(n_rows, n_cols, l, t, w, h).table

    # Basic styling
    for r in range(n_rows):
        for c in range(len(rows[r])): # specific row length
            cell = table.cell(r, c)
            cell.text = rows[r][c]
            # Optional: Add small font for tables so they fit
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12)

def assemble_slide(prs, slide_json, bg_path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    

    apply_background(prs, slide, bg_path)


    sorted_elements = sorted(slide_json["elements"], key=lambda x: x.get("z", 0))

    for el in sorted_elements:
        if el["type"] == "text":
            add_text(prs,slide, el)
        elif el["type"] == "image":
            add_image(prs,slide, el)
        elif el["type"] == "table":
            add_table(prs,slide, el)
            

